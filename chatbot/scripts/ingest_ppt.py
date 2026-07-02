#!/usr/bin/env python3
"""
PowerPoint Document Ingestion for RAG System

This script ingests a PowerPoint file into the RAG system:
1. Extracts text from slides
2. Chunks by slide (simple and effective)
3. Generates embeddings using Bedrock Titan
4. Stores in PostgreSQL with pgvector

Usage:
    python scripts/ingest_ppt.py --file methodology.pptx --env dev
"""

import argparse
import sys
import os
import json
import time
import boto3
import psycopg2
from pathlib import Path
from typing import List, Dict, Tuple
from pptx import Presentation

from env_config import EnvConfigError, exit_on_config_error, get_rds_connection_info

# Colors for output
GREEN = '\033[92m'
RED = '\033[91m'
YELLOW = '\033[93m'
BLUE = '\033[94m'
RESET = '\033[0m'


def print_header(text: str):
    """Print formatted header"""
    print(f"\n{BLUE}{'='*60}{RESET}")
    print(f"{BLUE}{text:^60}{RESET}")
    print(f"{BLUE}{'='*60}{RESET}\n")


def print_success(text: str):
    """Print success message"""
    print(f"{GREEN}✓{RESET} {text}")


def print_error(text: str):
    """Print error message"""
    print(f"{RED}✗{RESET} {text}")


def print_info(text: str):
    """Print info message"""
    print(f"{YELLOW}→{RESET} {text}")


def extract_text_from_ppt(file_path: str) -> List[Tuple[int, str, str]]:
    """
    Extract text from PowerPoint slides
    
    Returns:
        List of (slide_number, title, content) tuples
    """
    print_info(f"Extracting text from: {file_path}")
    
    prs = Presentation(file_path)
    slides_data = []
    
    for idx, slide in enumerate(prs.slides, start=1):
        # Extract title
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
        
        # Extract all text from shapes
        content_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text and text != title:  # Don't duplicate title
                    content_parts.append(text)
        
        content = "\n".join(content_parts)
        
        # Only include slides with content
        if title or content:
            slides_data.append((idx, title, content))
            print_success(f"Slide {idx}: {title[:50] if title else '(No title)'}...")
    
    print_success(f"Extracted {len(slides_data)} slides with content")
    return slides_data


def chunk_slides(slides_data: List[Tuple[int, str, str]]) -> List[Dict]:
    """
    Create chunks from slides
    
    For this demo, we use simple slide-based chunking:
    - 1 slide = 1 chunk (or 2 if slide is very long)
    
    Returns:
        List of chunk dictionaries
    """
    print_info("Creating chunks from slides...")
    
    chunks = []
    chunk_index = 0
    
    for slide_num, title, content in slides_data:
        # Combine title and content
        full_text = f"{title}\n\n{content}" if title else content
        
        # Simple chunking: if slide is short (<1000 chars), 1 chunk
        # If long, split into multiple chunks
        MAX_CHUNK_SIZE = 1000
        
        if len(full_text) <= MAX_CHUNK_SIZE:
            chunks.append({
                'chunk_index': chunk_index,
                'chunk_text': full_text,
                'slide_number': slide_num,
                'slide_title': title,
                'token_count': len(full_text.split())  # Rough estimate
            })
            chunk_index += 1
        else:
            # Split long slides into multiple chunks
            words = full_text.split()
            current_chunk = []
            current_size = 0
            
            for word in words:
                current_chunk.append(word)
                current_size += len(word) + 1
                
                if current_size >= MAX_CHUNK_SIZE:
                    chunks.append({
                        'chunk_index': chunk_index,
                        'chunk_text': ' '.join(current_chunk),
                        'slide_number': slide_num,
                        'slide_title': title,
                        'token_count': len(current_chunk)
                    })
                    chunk_index += 1
                    current_chunk = []
                    current_size = 0
            
            # Add remaining words
            if current_chunk:
                chunks.append({
                    'chunk_index': chunk_index,
                    'chunk_text': ' '.join(current_chunk),
                    'slide_number': slide_num,
                    'slide_title': title,
                    'token_count': len(current_chunk)
                })
                chunk_index += 1
    
    print_success(f"Created {len(chunks)} chunks")
    return chunks


def generate_embeddings(chunks: List[Dict], region: str) -> List[Dict]:
    """
    Generate embeddings for chunks using Bedrock Titan
    """
    print_info(f"Generating embeddings for {len(chunks)} chunks...")
    
    bedrock = boto3.client('bedrock-runtime', region_name=region)
    model_id = "amazon.titan-embed-text-v2:0"
    
    for i, chunk in enumerate(chunks, start=1):
        try:
            # Generate embedding
            body = json.dumps({
                "inputText": chunk['chunk_text'],
                "dimensions": 1024,
                "normalize": True
            })
            
            response = bedrock.invoke_model(
                modelId=model_id,
                body=body,
                contentType="application/json",
                accept="application/json"
            )
            
            result = json.loads(response['body'].read())
            chunk['embedding'] = result['embedding']
            
            print_success(f"Chunk {i}/{len(chunks)}: Generated 1024-dim embedding")
            
            # Small delay to avoid throttling
            if i < len(chunks):
                time.sleep(0.2)
        
        except Exception as e:
            print_error(f"Failed to generate embedding for chunk {i}: {str(e)}")
            raise
    
    print_success("All embeddings generated successfully")
    return chunks


def get_db_credentials(secret_arn: str, region: str) -> Dict[str, str]:
    """Retrieve database credentials from Secrets Manager"""
    client = boto3.client('secretsmanager', region_name=region)
    response = client.get_secret_value(SecretId=secret_arn)
    return json.loads(response['SecretString'])


def get_db_connection_info(env_name: str, region: str) -> Dict[str, str]:
    """Get database connection details from environment YAML."""
    del region
    try:
        return get_rds_connection_info(env_name)
    except EnvConfigError as exc:
        exit_on_config_error(exc)


def upload_to_s3(file_path: str, env_name: str, region: str) -> Tuple[str, str]:
    """Upload file to S3 and return bucket name and key"""
    print_info("Uploading file to S3...")
    
    # Get S3 bucket name from CloudFormation
    cf_client = boto3.client('cloudformation', region_name=region)
    stack_name = f"ChatbotS3Stack-{env_name}"
    
    response = cf_client.describe_stacks(StackName=stack_name)
    outputs = response['Stacks'][0]['Outputs']
    bucket_name = None
    
    for output in outputs:
        if output['OutputKey'] == 'DocumentsBucketName':
            bucket_name = output['OutputValue']
            break
    
    if not bucket_name:
        print_error("Could not find S3 bucket name")
        sys.exit(1)
    
    # Upload file
    s3_client = boto3.client('s3', region_name=region)
    file_name = Path(file_path).name
    s3_key = f"documents/{file_name}"
    
    with open(file_path, 'rb') as f:
        s3_client.put_object(
            Bucket=bucket_name,
            Key=s3_key,
            Body=f,
            ServerSideEncryption='AES256'
        )
    
    print_success(f"Uploaded to s3://{bucket_name}/{s3_key}")
    return bucket_name, s3_key


def store_in_database(
    file_path: str,
    bucket_name: str,
    s3_key: str,
    chunks: List[Dict],
    db_info: Dict[str, str],
    region: str
):
    """Store document and chunks in PostgreSQL"""
    print_info("Storing in database...")
    
    # Get credentials
    credentials = get_db_credentials(db_info['secret_arn'], region)
    
    # Connect to database
    conn = psycopg2.connect(
        host=db_info['host'],
        port=int(db_info['port']),
        dbname=db_info['dbname'],
        user=credentials['username'],
        password=credentials['password'],
        connect_timeout=30
    )
    cursor = conn.cursor()
    
    try:
        # Insert document metadata
        file_name = Path(file_path).name
        file_size = os.path.getsize(file_path)
        
        cursor.execute("""
            INSERT INTO documents (
                filename, s3_key, s3_bucket, file_type, 
                file_size_bytes, status, metadata
            )
            VALUES (%s, %s, %s, %s, %s, %s, %s)
            RETURNING id;
        """, (
            file_name,
            s3_key,
            bucket_name,
            'pptx',
            file_size,
            'ready',
            json.dumps({'total_chunks': len(chunks)})
        ))
        
        document_id = cursor.fetchone()[0]
        print_success(f"Document stored (ID: {document_id})")
        
        # Insert chunks with embeddings
        for chunk in chunks:
            cursor.execute("""
                INSERT INTO document_chunks (
                    document_id, chunk_index, chunk_text,
                    chunk_embedding, token_count
                )
                VALUES (%s, %s, %s, %s::vector, %s);
            """, (
                document_id,
                chunk['chunk_index'],
                chunk['chunk_text'],
                chunk['embedding'],
                chunk['token_count']
            ))
        
        conn.commit()
        print_success(f"Stored {len(chunks)} chunks with embeddings")
        
        # Verify
        cursor.execute("""
            SELECT COUNT(*) FROM document_chunks WHERE document_id = %s;
        """, (document_id,))
        count = cursor.fetchone()[0]
        print_success(f"Verified: {count} chunks in database")
        
    except Exception as e:
        conn.rollback()
        print_error(f"Database error: {str(e)}")
        raise
    finally:
        cursor.close()
        conn.close()


def main():
    parser = argparse.ArgumentParser(
        description="Ingest PowerPoint document into RAG system"
    )
    parser.add_argument(
        "--file",
        type=str,
        required=True,
        help="Path to PowerPoint file"
    )
    parser.add_argument(
        "--env",
        type=str,
        default="stage",
        choices=["dev", "stage", "uat", "prod"],
        help="Environment name"
    )
    parser.add_argument(
        "--region",
        type=str,
        default="us-east-1",
        help="AWS region"
    )
    
    args = parser.parse_args()
    
    # Validate file exists
    if not os.path.exists(args.file):
        print_error(f"File not found: {args.file}")
        sys.exit(1)
    
    if not args.file.endswith('.pptx'):
        print_error("File must be a PowerPoint (.pptx) file")
        sys.exit(1)
    
    print_header("PPT DOCUMENT INGESTION")
    print(f"File: {args.file}")
    print(f"Environment: {args.env}")
    print(f"Region: {args.region}\n")
    
    try:
        # Step 1: Extract text from PPT
        print_header("STEP 1: Extract Text")
        slides_data = extract_text_from_ppt(args.file)
        
        # Step 2: Create chunks
        print_header("STEP 2: Create Chunks")
        chunks = chunk_slides(slides_data)
        
        # Step 3: Generate embeddings
        print_header("STEP 3: Generate Embeddings")
        chunks = generate_embeddings(chunks, args.region)
        
        # Step 4: Upload to S3
        print_header("STEP 4: Upload to S3")
        bucket_name, s3_key = upload_to_s3(args.file, args.env, args.region)
        
        # Step 5: Store in database
        print_header("STEP 5: Store in Database")
        db_info = get_db_connection_info(args.env, args.region)
        
        # Note: This will fail if run from outside VPC
        # User should use the ingestion Lambda instead (TODO: create that)
        # For now, document this limitation
        print_info("Note: Database connection requires VPC access")
        print_info("If this fails, the Lambda ingestion function is needed")
        
        store_in_database(
            args.file,
            bucket_name,
            s3_key,
            chunks,
            db_info,
            args.region
        )
        
        # Success!
        print_header("INGESTION COMPLETE")
        print_success(f"Document: {Path(args.file).name}")
        print_success(f"Total slides: {len(slides_data)}")
        print_success(f"Total chunks: {len(chunks)}")
        print_success(f"S3 location: s3://{bucket_name}/{s3_key}")
        print_success("Database: Stored with embeddings")
        print()
        print("Next steps:")
        print("1. Deploy updated runtime Lambda with RAG endpoint")
        print("2. Test with: POST /v1/chat-rag")
        print()
        
    except Exception as e:
        print_error(f"Ingestion failed: {str(e)}")
        sys.exit(1)


if __name__ == "__main__":
    main()
