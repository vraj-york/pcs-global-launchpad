"""
Document Ingestion Lambda for RAG System

This Lambda ingests documents from S3 into the RAG database.
Refactored to use shared utilities for consistency.

Invocation:
    aws lambda invoke --function-name bispy-bot-ingest-dev \
        --payload '{"s3_bucket": "bucket", "s3_key": "key"}' response.json
"""

import json
import os
import logging
import boto3
from typing import List, Tuple
from io import BytesIO
from pptx import Presentation
from docx import Document

# Import from shared package
from shared.chunking import chunk_sections
from shared.embeddings import generate_embeddings_batch
from shared.database_utils import get_db_connection

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)


def extract_text_from_ppt(ppt_bytes: bytes) -> List[Tuple[int, str, str]]:
    """
    Extract text from PowerPoint presentation

    Args:
        ppt_bytes: PowerPoint file bytes

    Returns:
        List of (slide_number, title, content) tuples
    """
    prs = Presentation(BytesIO(ppt_bytes))
    slides_data = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        content_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text and text != title:
                    content_parts.append(text)

        content = "\n".join(content_parts)

        if title or content:
            slides_data.append((idx, title, content))

    return slides_data


def extract_text_from_docx(docx_bytes: bytes) -> List[Tuple[int, str, str]]:
    """
    Extract text from Word document

    Args:
        docx_bytes: Word document file bytes

    Returns:
        List of (section_number, heading, content) tuples
    """
    doc = Document(BytesIO(docx_bytes))
    sections_data = []

    current_section = 1
    current_heading = ""
    current_content = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Check if this is a heading
        if para.style.name.startswith("Heading"):
            # Save previous section
            if current_heading or current_content:
                content = "\n".join(current_content)
                if current_heading or content:
                    sections_data.append((current_section, current_heading, content))
                    current_section += 1

            # Start new section
            current_heading = text
            current_content = []
        else:
            current_content.append(text)

    # Add final section
    if current_heading or current_content:
        content = "\n".join(current_content)
        if current_heading or content:
            sections_data.append((current_section, current_heading, content))

    # Fallback: treat entire doc as one section
    if not sections_data and doc.paragraphs:
        all_text = "\n".join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
        if all_text:
            sections_data.append((1, "Document Content", all_text))

    return sections_data


def handler(event, context):
    """
    Ingest document from S3 into RAG database

    Event format:
    {
        "s3_bucket": "bucket-name",
        "s3_key": "path/to/file.pptx",
        "allowed_roles": ["superadmin", "coach", "employee", "default"]  # Optional
    }
    
    If allowed_roles is not specified, defaults to ["superadmin"] (admin-only access).
    This ensures secure-by-default behavior where new documents are restricted unless
    explicitly made accessible to other roles.
    """
    import asyncio
    
    # Run async handler in event loop
    return asyncio.run(async_handler(event, context))


async def async_handler(event, context):
    """
    Async version of the handler for non-blocking embedding generation
    """
    logger.info(f"Ingestion Lambda invoked: {json.dumps(event)}")

    try:
        # 1. Get S3 location and roles
        missing = [f for f in ("s3_bucket", "s3_key") if f not in event]
        if missing:
            raise ValueError(
                f"Missing required field(s) in event payload: {missing}. "
                "Expected: {\"s3_bucket\": \"<bucket-name>\", \"s3_key\": \"<bucket-file-path>\"}"
            )

        s3_bucket = event["s3_bucket"]
        s3_key = event["s3_key"]
        
        # Get allowed_roles from event, default to superadmin-only if not specified
        # This ensures secure-by-default behavior
        allowed_roles = event.get("allowed_roles", ["superadmin"])
        
        logger.info(f"Processing: s3://{s3_bucket}/{s3_key}")
        logger.info(f"Allowed roles: {allowed_roles}")

        # 2. Download file from S3
        s3_client = boto3.client("s3")
        response = s3_client.get_object(Bucket=s3_bucket, Key=s3_key)
        file_bytes = response["Body"].read()
        file_size = len(file_bytes)
        logger.info(f"Downloaded {file_size} bytes")

        # 3. Extract text based on file type
        file_extension = s3_key.lower().split(".")[-1]

        if file_extension == "pptx":
            logger.info("Extracting from PowerPoint...")
            sections = extract_text_from_ppt(file_bytes)
            section_type = "slide"
        elif file_extension == "docx":
            logger.info("Extracting from Word document...")
            sections = extract_text_from_docx(file_bytes)
            section_type = "section"
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")

        logger.info(f"Extracted {len(sections)} {section_type}s")

        # 4. Create chunks using shared utility with config from environment
        chunk_size = int(os.environ.get("CHUNK_SIZE", "2000"))
        chunk_overlap = int(os.environ.get("CHUNK_OVERLAP", "200"))
        logger.info(f"Creating chunks with size={chunk_size}, overlap={chunk_overlap}...")
        chunks = chunk_sections(sections, chunk_size=chunk_size, section_type=section_type)
        logger.info(f"Created {len(chunks)} chunks")

        # 5. Generate embeddings using shared utility (async)
        logger.info("Generating embeddings (async)...")
        chunk_texts = [chunk["chunk_text"] for chunk in chunks]
        embeddings = await generate_embeddings_batch(chunk_texts)

        # Attach embeddings to chunks
        for chunk, embedding in zip(chunks, embeddings):
            chunk["embedding"] = embedding

        logger.info(f"Generated {len(embeddings)} embeddings")

        # 6. Store in database
        logger.info("Storing in database...")
        db_host = os.environ["DB_HOST"]
        db_port = int(os.environ["DB_PORT"])
        db_name = os.environ["DB_NAME"]
        secret_arn = os.environ["DB_SECRET_ARN"]

        with get_db_connection(db_host, db_port, db_name, secret_arn) as conn:
            cursor = conn.cursor()

            # Insert document metadata with role-based access control
            filename = s3_key.split("/")[-1]
            cursor.execute(
                """
                INSERT INTO documents (
                    filename, s3_key, s3_bucket, file_type,
                    file_size_bytes, status, metadata, allowed_roles
                )
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s::jsonb)
                RETURNING id;
            """,
                (
                    filename,
                    s3_key,
                    s3_bucket,
                    file_extension,
                    file_size,
                    "ready",
                    json.dumps({"total_chunks": len(chunks)}),
                    json.dumps(allowed_roles),  # Convert list to JSON for JSONB column
                ),
            )

            document_id = cursor.fetchone()[0]

            # Insert chunks with embeddings
            for chunk in chunks:
                cursor.execute(
                    """
                    INSERT INTO document_chunks (
                        document_id, chunk_index, chunk_text,
                        chunk_embedding, token_count
                    )
                    VALUES (%s, %s, %s, %s::vector, %s);
                """,
                    (
                        document_id,
                        chunk["chunk_index"],
                        chunk["chunk_text"],
                        chunk["embedding"],
                        chunk["token_count"],
                    ),
                )

            conn.commit()

        logger.info(f"Stored document (ID: {document_id}) with roles: {allowed_roles}")

        return {
            "statusCode": 200,
            "body": json.dumps(
                {
                    "status": "success",
                    "message": "Document ingested successfully",
                    "document_id": document_id,
                    "file_type": file_extension,
                    "total_sections": len(sections),
                    "total_chunks": len(chunks),
                    "s3_location": f"s3://{s3_bucket}/{s3_key}",
                    "allowed_roles": allowed_roles,
                }
            ),
        }

    except Exception as e:
        error_msg = f"Ingestion failed: {str(e)}"
        logger.error(f"ERROR: {error_msg}", exc_info=True)

        return {
            "statusCode": 500,
            "body": json.dumps({"status": "error", "message": error_msg}),
        }
